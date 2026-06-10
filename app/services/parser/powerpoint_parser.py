from __future__ import annotations

import re
import time
from collections import Counter
from dataclasses import dataclass
from io import BytesIO
from statistics import median
from pathlib import Path
from typing import Any, Iterable

from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    from app.services.parser.img_parser import enrich_image_descriptions
    from app.services.parser.word_parser import format_extracted_items
except ModuleNotFoundError:
    from img_parser import enrich_image_descriptions
    from word_parser import format_extracted_items


IMAGE_EXTENSIONS = {
    "image/png": ".png",
    "image/jpeg": ".jpg",
    "image/jpg": ".jpg",
    "image/gif": ".gif",
    "image/bmp": ".bmp",
    "image/tiff": ".tiff",
    "image/x-emf": ".emf",
    "image/x-wmf": ".wmf",
}
TITLE_KEYWORD_PATTERN = re.compile(
    r"(?i)\b(part|chapter|section)\b|章节|单元|模块|(?:^|[^\d])(?:0?[1-9]|[1-9]\d)(?:[^\d]|$)"
)
CATALOG_KEYWORDS = ("目录", "目 录", "contents", "agenda")
MAX_TITLE_PAGE_CHARS = 20
HASH_SIZE = 8
REPEATED_HASH_DISTANCE = 8
POSITION_TOLERANCE = 0.04
SIZE_TOLERANCE = 0.06
ROW_TOP_TOLERANCE_EMU = 260_000
FRONT_MATTER_SLIDE_LIMIT = 5
TITLE_STYLE = "标题 1"
PAGE_HEADING_STYLE = "标题 2"


@dataclass
class ParseContext:
    source_path: Path
    image_dir: Path
    image_index: int = 0

    def next_image_path(self, content_type: str | None) -> Path:
        self.image_index += 1
        extension = IMAGE_EXTENSIONS.get(str(content_type or "").lower(), ".bin")
        return self.image_dir / f"image_{self.image_index:04d}{extension}"


def parse_powerpoint_document(
    file_path: str | Path,
    *,
    image_analysis_workers: int = 3,
) -> list[dict[str, Any]]:
    """Extract pptx slide text and images, remove repeated visual chrome, then describe images."""
    started_at = time.perf_counter()
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"PowerPoint document not found: {path}")
    if path.suffix.lower() != ".pptx":
        raise ValueError(f"Only .pptx files are supported: {path}")
    if path.name.startswith("~$"):
        raise ValueError(f"Temporary PowerPoint lock files are not supported: {path}")

    _log(f"start parsing: {path}")
    presentation = Presentation(path)
    image_dir = Path("data") / "processing" / path.stem / "img"
    image_dir.mkdir(parents=True, exist_ok=True)
    _remove_stale_bin_images(image_dir)

    stage_started_at = time.perf_counter()
    context = ParseContext(source_path=path, image_dir=image_dir)
    slides = _extract_slides(presentation, context)
    _log(f"extracted slides: parsed={len(slides)}, images={context.image_index} ({time.perf_counter() - stage_started_at:.2f}s)")

    stage_started_at = time.perf_counter()
    items = _build_items(slides)
    _log_item_summary("built slide items", items, stage_started_at)

    stage_started_at = time.perf_counter()
    items = _remove_repeated_images(items, slide_count=len(slides))
    _log_item_summary("removed repeated images", items, stage_started_at)

    stage_started_at = time.perf_counter()
    _log("start image analysis")
    items = enrich_image_descriptions(items, document_title=path.stem, max_concurrency=image_analysis_workers)
    _log_item_summary("finished image analysis", items, stage_started_at)

    stage_started_at = time.perf_counter()
    output_path = write_items_to_txt(items, path)
    _log(f"wrote parsed txt: {output_path} ({time.perf_counter() - stage_started_at:.2f}s)")
    _log(f"finished parsing: {path.name} ({time.perf_counter() - started_at:.2f}s)")
    return items


def _extract_slides(presentation: Presentation, context: ParseContext) -> list[dict[str, Any]]:
    output: list[dict[str, Any]] = []
    last_slide_index = len(presentation.slides)

    for slide_index, slide in enumerate(presentation.slides, start=1):
        if slide_index == last_slide_index:
            _log(f"skip last slide: {slide_index}")
            continue
        if _slide_has_video(slide):
            _log(f"skip video slide: {slide_index}")
            continue

        slide_width = int(presentation.slide_width)
        slide_height = int(presentation.slide_height)
        elements: list[dict[str, Any]] = []

        for shape_order, shape in enumerate(_iter_shapes(slide.shapes), start=1):
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                image_item = _extract_picture(shape, context, slide_index, shape_order, slide_width, slide_height)
                if image_item is not None:
                    elements.append(image_item)
                continue

            for text_item in _extract_text_items(shape, slide_index, shape_order):
                elements.append(text_item)

        elements = _sort_elements_reading_order(elements)
        output.append({"slide_index": slide_index, "elements": elements})

    return output


def _iter_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _slide_has_video(slide: Any) -> bool:
    for shape in _iter_shapes(slide.shapes):
        shape_type = getattr(shape, "shape_type", None)
        if shape_type in {MSO_SHAPE_TYPE.MEDIA, MSO_SHAPE_TYPE.WEB_VIDEO}:
            return True
    return False


def _extract_picture(
    shape: Any,
    context: ParseContext,
    slide_index: int,
    shape_order: int,
    slide_width: int,
    slide_height: int,
) -> dict[str, Any] | None:
    image = getattr(shape, "image", None)
    if image is None:
        return None

    image_path = context.next_image_path(getattr(image, "content_type", None))
    blob = image.blob
    image_path.write_bytes(blob)
    left = int(getattr(shape, "left", 0) or 0)
    top = int(getattr(shape, "top", 0) or 0)
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)

    return {
        "type": "image",
        "style": "图片",
        "source": f"slide:{slide_index}:shape:{shape_order}",
        "text": str(image_path),
        "path": str(image_path),
        "slide": str(slide_index),
        "left": str(left),
        "top": str(top),
        "width": str(width),
        "height": str(height),
        "size": str(len(blob)),
        "norm_left": left / slide_width if slide_width else 0.0,
        "norm_top": top / slide_height if slide_height else 0.0,
        "norm_width": width / slide_width if slide_width else 0.0,
        "norm_height": height / slide_height if slide_height else 0.0,
        "order": shape_order,
        "hash": _image_dhash(blob),
    }


def _extract_text_items(shape: Any, slide_index: int, shape_order: int) -> list[dict[str, Any]]:
    if getattr(shape, "has_table", False):
        return _extract_table_text_items(shape, slide_index, shape_order)
    if not getattr(shape, "has_text_frame", False):
        return []

    text = _shape_text(shape)
    if not text:
        return []

    return [
        {
            "type": "paragraph",
            "style": "正文",
            "text": text,
            "source": f"slide:{slide_index}:shape:{shape_order}",
            "slide": str(slide_index),
            "font_size": _shape_max_font_size(shape),
            "left": str(int(getattr(shape, "left", 0) or 0)),
            "top": str(int(getattr(shape, "top", 0) or 0)),
            "height": str(int(getattr(shape, "height", 0) or 0)),
            "order": shape_order,
        }
    ]


def _extract_table_text_items(shape: Any, slide_index: int, shape_order: int) -> list[dict[str, Any]]:
    rows: list[list[str]] = []
    for row in shape.table.rows:
        values = [_clean_text(cell.text) for cell in row.cells]
        if any(values):
            rows.append(values)
    if not rows:
        return []
    text = "\n".join(" | ".join(value for value in row if value) for row in rows)
    return [
        {
            "type": "paragraph",
            "style": "正文",
            "text": text,
            "source": f"slide:{slide_index}:shape:{shape_order}:table",
            "slide": str(slide_index),
            "font_size": _shape_max_font_size(shape),
            "left": str(int(getattr(shape, "left", 0) or 0)),
            "top": str(int(getattr(shape, "top", 0) or 0)),
            "height": str(int(getattr(shape, "height", 0) or 0)),
            "order": shape_order,
        }
    ]


def _shape_text(shape: Any) -> str:
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        return ""
    lines = [_clean_text(paragraph.text) for paragraph in text_frame.paragraphs]
    return "\n".join(line for line in lines if line).strip()


def _shape_max_font_size(shape: Any) -> float:
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        return 0.0
    max_size = 0.0
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                max_size = max(max_size, float(run.font.size.pt))
        if paragraph.font.size is not None:
            max_size = max(max_size, float(paragraph.font.size.pt))
    return max_size


def _build_items(slides: list[dict[str, Any]]) -> list[dict[str, Any]]:
    slide_infos = _slide_infos(slides)
    title_cutoff = _title_slide_font_cutoff([float(info["max_font_size"]) for info in slide_infos])
    title_slide_indexes = {
        int(info["slide_index"])
        for info in slide_infos
        if _looks_like_title_slide(info, title_cutoff)
    }
    catalog_slide_indexes = {
        int(info["slide_index"])
        for info in slide_infos
        if _looks_like_catalog_slide(info, title_cutoff)
    }
    output: list[dict[str, Any]] = []
    page_number = 1

    for slide in slides:
        slide_index = int(slide["slide_index"])
        if slide_index in catalog_slide_indexes:
            _log(f"skip catalog slide: {slide_index}")
            continue

        if slide_index in title_slide_indexes:
            title_item = _merged_title_item(slide)
            if title_item is not None:
                output.append(title_item)
            page_number = 1
            continue

        text_elements = [element for element in slide["elements"] if element.get("type") == "paragraph"]
        if text_elements:
            output.append(
                {
                    "type": "paragraph",
                    "style": PAGE_HEADING_STYLE,
                    "text": f"第{page_number}页",
                    "source": f"slide:{slide_index}:page",
                    "slide": str(slide_index),
                }
            )
            page_number += 1

        for element in slide["elements"]:
            if slide_index == 1 and element.get("type") == "image":
                continue
            item = dict(element)
            output.append(item)

    return output


def _slide_infos(slides: list[dict[str, Any]]) -> list[dict[str, Any]]:
    return [
        info
        for info in (_title_slide_info(slide) for slide in slides)
        if info["text"]
    ]


def _title_slide_info(slide: dict[str, Any]) -> dict[str, Any]:
    text_items = [
        element
        for element in slide["elements"]
        if element.get("type") == "paragraph" and str(element.get("text") or "").strip()
    ]
    text = "\n".join(str(item.get("text") or "").strip() for item in text_items)
    max_font_size = max((float(item.get("font_size") or 0) for item in text_items), default=0.0)
    return {
        "slide_index": int(slide["slide_index"]),
        "text": text,
        "compact_text": re.sub(r"\s+", "", text),
        "max_font_size": max_font_size,
    }


def _title_slide_font_cutoff(sizes: list[float]) -> float:
    sizes = sorted(size for size in sizes if size > 0)
    if not sizes:
        return 24.0
    median_size = float(median(sizes))
    percentile_75 = sizes[int((len(sizes) - 1) * 0.75)]
    return max(24.0, min(percentile_75, median_size + 4.0))


def _looks_like_title_slide(info: dict[str, Any], font_cutoff: float) -> bool:
    text = str(info.get("text") or "").strip()
    compact = str(info.get("compact_text") or "")
    if not compact or len(compact) >= MAX_TITLE_PAGE_CHARS:
        return False
    if any(keyword.lower() in text.lower() for keyword in CATALOG_KEYWORDS):
        return False
    if float(info.get("max_font_size") or 0) < font_cutoff:
        return False
    return bool(TITLE_KEYWORD_PATTERN.search(text))


def _looks_like_catalog_slide(info: dict[str, Any], font_cutoff: float) -> bool:
    slide_index = int(info.get("slide_index") or 0)
    if slide_index > FRONT_MATTER_SLIDE_LIMIT:
        return False
    text = str(info.get("text") or "").strip()
    if not any(keyword.lower() in text.lower() for keyword in CATALOG_KEYWORDS):
        return False
    if float(info.get("max_font_size") or 0) < font_cutoff:
        return False
    return bool(TITLE_KEYWORD_PATTERN.search(text))


def _merged_title_item(slide: dict[str, Any]) -> dict[str, Any] | None:
    text = _merged_slide_text(slide)
    if not text:
        return None
    slide_index = int(slide["slide_index"])
    return {
        "type": "paragraph",
        "style": TITLE_STYLE,
        "text": text,
        "source": f"slide:{slide_index}:title",
        "slide": str(slide_index),
    }


def _merged_slide_text(slide: dict[str, Any]) -> str:
    parts = [
        str(element.get("text") or "").strip()
        for element in slide["elements"]
        if element.get("type") == "paragraph" and str(element.get("text") or "").strip()
    ]
    return " ".join(parts).strip()


def _remove_repeated_images(items: list[dict[str, Any]], *, slide_count: int) -> list[dict[str, Any]]:
    image_items = [item for item in items if item.get("type") == "image" and item.get("hash")]
    if len(image_items) < 2:
        return items

    clusters: list[list[dict[str, Any]]] = []
    for item in image_items:
        for cluster in clusters:
            if _same_repeated_image_group(item, cluster[0]):
                cluster.append(item)
                break
        else:
            clusters.append([item])

    min_cluster_size = 2
    removed_paths: set[str] = set()
    for cluster in clusters:
        if len(cluster) < min_cluster_size:
            continue
        if _max_hash_distance(cluster) <= REPEATED_HASH_DISTANCE:
            removed_paths.update(str(item.get("path") or "") for item in cluster)

    if not removed_paths:
        return items
    _log(f"removed repeated image clusters: images={len(removed_paths)}")
    return [
        item
        for item in items
        if item.get("type") != "image" or str(item.get("path") or "") not in removed_paths
    ]


def _same_repeated_image_group(left: dict[str, Any], right: dict[str, Any]) -> bool:
    if _hash_distance(str(left.get("hash")), str(right.get("hash"))) > REPEATED_HASH_DISTANCE:
        return False
    return (
        abs(float(left.get("norm_left") or 0) - float(right.get("norm_left") or 0)) <= POSITION_TOLERANCE
        and abs(float(left.get("norm_top") or 0) - float(right.get("norm_top") or 0)) <= POSITION_TOLERANCE
        and abs(float(left.get("norm_width") or 0) - float(right.get("norm_width") or 0)) <= SIZE_TOLERANCE
        and abs(float(left.get("norm_height") or 0) - float(right.get("norm_height") or 0)) <= SIZE_TOLERANCE
    )


def _max_hash_distance(cluster: list[dict[str, Any]]) -> int:
    max_distance = 0
    hashes = [str(item.get("hash")) for item in cluster]
    for left_index, left_hash in enumerate(hashes):
        for right_hash in hashes[left_index + 1 :]:
            max_distance = max(max_distance, _hash_distance(left_hash, right_hash))
    return max_distance


def _image_dhash(blob: bytes) -> str:
    try:
        with Image.open(BytesIO(blob)) as image:
            grayscale = image.convert("L").resize((HASH_SIZE + 1, HASH_SIZE), Image.Resampling.LANCZOS)
    except (UnidentifiedImageError, OSError):
        return ""

    pixels = list(grayscale.getdata())
    bits: list[str] = []
    for row in range(HASH_SIZE):
        offset = row * (HASH_SIZE + 1)
        for column in range(HASH_SIZE):
            bits.append("1" if pixels[offset + column] > pixels[offset + column + 1] else "0")
    return f"{int(''.join(bits), 2):016x}"


def _hash_distance(left_hash: str, right_hash: str) -> int:
    if not left_hash or not right_hash:
        return 64
    try:
        return (int(left_hash, 16) ^ int(right_hash, 16)).bit_count()
    except ValueError:
        return 64


def _sort_elements_reading_order(elements: list[dict[str, Any]]) -> list[dict[str, Any]]:
    rows: list[list[dict[str, Any]]] = []

    for item in sorted(elements, key=_vertical_sort_key):
        top = int(item.get("top") or 0)
        for row in rows:
            row_top = _row_top(row)
            if abs(top - row_top) <= ROW_TOP_TOLERANCE_EMU:
                row.append(item)
                break
        else:
            rows.append([item])

    output: list[dict[str, Any]] = []
    for row in sorted(rows, key=_row_top):
        output.extend(sorted(row, key=_horizontal_sort_key))
    return output


def _vertical_sort_key(item: dict[str, Any]) -> tuple[int, int, int]:
    return (
        int(item.get("top") or 0),
        int(item.get("left") or 0),
        int(item.get("order") or 0),
    )


def _horizontal_sort_key(item: dict[str, Any]) -> tuple[int, int, int]:
    return (
        int(item.get("left") or 0),
        int(item.get("top") or 0),
        int(item.get("order") or 0),
    )


def _row_top(row: list[dict[str, Any]]) -> int:
    if not row:
        return 0
    return min(int(item.get("top") or 0) for item in row)


def _clean_text(text: str) -> str:
    lines = [re.sub(r"[ \t\u3000]+", " ", line).strip() for line in str(text or "").splitlines()]
    return "\n".join(line for line in lines if line).strip()


def write_items_to_txt(items: list[dict[str, Any]], source_file: str | Path) -> Path:
    source_path = Path(source_file)
    txt_dir = Path("data") / "processing" / source_path.stem / "txt"
    txt_dir.mkdir(parents=True, exist_ok=True)

    output_path = txt_dir / f"{source_path.stem}.txt"
    output_path.write_text(format_extracted_items(items), encoding="utf-8")
    return output_path


def _remove_stale_bin_images(image_dir: Path) -> None:
    for path in image_dir.glob("*.bin"):
        path.unlink(missing_ok=True)


def _log(message: str) -> None:
    print(f"[powerpoint_parser] {message}", flush=True)


def _log_item_summary(message: str, items: list[dict[str, Any]], started_at: float) -> None:
    counts = Counter(item.get("type", "unknown") for item in items)
    type_summary = ", ".join(f"{item_type}={count}" for item_type, count in sorted(counts.items()))
    _log(f"{message}: total={len(items)} ({type_summary or 'no items'}) ({time.perf_counter() - started_at:.2f}s)")


if __name__ == "__main__":
    target_file = Path("data") / "raw" / "clear_theme_pptx" / "sample.pptx"
    extracted_items = parse_powerpoint_document(target_file)
    txt_path = Path("data") / "processing" / target_file.stem / "txt" / f"{target_file.stem}.txt"

    print(f"File: {target_file}")
    print(f"Extracted text blocks: {len(extracted_items)}")
    print(f"Text output: {txt_path}")
